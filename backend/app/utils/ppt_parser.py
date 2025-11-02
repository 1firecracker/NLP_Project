"""PPT 文档解析器"""
from pptx import Presentation
from typing import List, Dict, Any
from pathlib import Path


class PPTParser:
    """PPTX 文档解析器"""
    
    def parse(self, file_path: str) -> Dict[str, Any]:
        """解析 PPTX 文件，返回结构化数据
        
        Args:
            file_path: PPTX 文件路径
            
        Returns:
            包含幻灯片信息的字典
        """
        prs = Presentation(file_path)
        
        slides_data = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": idx,
                "title": self._extract_title(slide),
                "text_content": self._extract_text(slide),
                "images": self._extract_images(slide, idx),
                "structure": self._extract_structure(slide),
            }
            slides_data.append(slide_data)
        
        return {
            "filename": Path(file_path).name,
            "total_slides": len(slides_data),
            "slides": slides_data
        }
    
    def extract_text(self, file_path: str) -> str:
        """提取 PPT 的纯文本内容
        
        Args:
            file_path: PPTX 文件路径
            
        Returns:
            所有幻灯片文本内容（用换行符分隔）
        """
        prs = Presentation(file_path)
        texts = []
        
        for slide in prs.slides:
            slide_texts = []
            # 提取标题
            title = self._extract_title(slide)
            if title:
                slide_texts.append(title)
            # 提取内容
            content = self._extract_text(slide)
            if content:
                slide_texts.append(content)
            
            if slide_texts:
                texts.append("\n".join(slide_texts))
        
        return "\n\n".join(texts)
    
    def _extract_title(self, slide) -> str:
        """提取幻灯片标题"""
        # 尝试从占位符获取标题
        for shape in slide.shapes:
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                if shape.placeholder_format.idx == 0:  # 标题占位符通常是 idx=0
                    if hasattr(shape, "text") and shape.text:
                        return shape.text.strip()
        
        # 如果没有找到占位符标题，尝试查找第一个文本框（通常是标题）
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text and len(text) < 200:  # 标题通常较短
                    return text
        
        return ""
    
    def _extract_text(self, slide) -> str:
        """提取幻灯片文本内容"""
        texts = []
        title_found = False
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    # 跳过标题（已在标题中提取）
                    if not title_found and len(text) < 200:
                        title_found = True
                        # 如果已经在标题中，跳过
                        continue
                    texts.append(text)
        
        return "\n".join(texts)
    
    def _extract_images(self, slide, slide_number: int) -> List[Dict[str, Any]]:
        """提取幻灯片中的图片元信息
        
        Args:
            slide: 幻灯片对象
            slide_number: 幻灯片编号
            
        Returns:
            图片信息列表
        """
        images = []
        for idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "image"):
                image_info = {
                    "image_id": f"slide_{slide_number}_img_{idx}",
                    "position": {
                        "left": int(shape.left),
                        "top": int(shape.top),
                        "width": int(shape.width),
                        "height": int(shape.height),
                    },
                    "alt_text": getattr(shape, "name", f"Image {idx + 1}")
                }
                images.append(image_info)
        return images
    
    def _extract_structure(self, slide) -> Dict[str, Any]:
        """提取幻灯片结构信息"""
        return {
            "layout": getattr(slide.slide_layout, "name", "Unknown"),
            "shapes_count": len(slide.shapes)
        }
    
    def extract_text_positions(self, file_path: str, slide_number: int) -> Dict[str, Any]:
        """提取指定幻灯片的文本及其位置信息（用于Canvas精确叠加）
        
        Args:
            file_path: PPTX文件路径
            slide_number: 幻灯片编号（从1开始）
            
        Returns:
            包含文本位置列表和幻灯片尺寸信息的字典
        """
        prs = Presentation(file_path)
        
        if slide_number < 1 or slide_number > len(prs.slides):
            return {
                "text_items": [],
                "slide_dimensions": None
            }
        
        slide = prs.slides[slide_number - 1]
        text_items = []
        
        # ✅ 获取实际幻灯片尺寸（不再硬编码）
        slide_width_emu = prs.slide_width   # 实际宽度（EMU单位）
        slide_height_emu = prs.slide_height  # 实际高度（EMU单位）
        
        # 转换为英寸
        # 1 inch = 914400 EMU
        slide_width_inches = slide_width_emu / 914400
        slide_height_inches = slide_height_emu / 914400
        
        # 基于150 DPI转换为像素（文本位置的基准DPI）
        resolution = 150
        slide_width_px = slide_width_inches * resolution
        slide_height_px = slide_height_inches * resolution
        
        # 缩放因子：像素/EMU（基于150 DPI）
        scale_x = slide_width_px / slide_width_emu
        scale_y = slide_height_px / slide_height_emu
        
        for shape in slide.shapes:
            # 尝试多种方式提取文本
            text = None
            
            # 方法1: 直接访问text属性
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
            
            # 方法2: 通过text_frame提取（更可靠）
            if not text and hasattr(shape, "text_frame"):
                try:
                    text_parts = []
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                text_parts.append(run.text)
                    if text_parts:
                        text = " ".join(text_parts).strip()
                except Exception as e:
                    # 忽略提取错误，继续尝试其他方法
                    pass
            
            # 如果找到文本，提取位置信息
            if text and len(text) > 0:
                try:
                    # 获取形状位置（EMU单位）
                    left_emu = int(shape.left)
                    top_emu = int(shape.top)
                    width_emu = int(shape.width)
                    height_emu = int(shape.height)
                    
                    # 转换为像素
                    x = left_emu * scale_x
                    y = top_emu * scale_y
                    width = width_emu * scale_x
                    height = height_emu * scale_y
                    
                    # 尝试获取字体大小
                    font_size = 16  # 默认值
                    try:
                        if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                            para = shape.text_frame.paragraphs[0]
                            if para.runs:
                                font = para.runs[0].font
                                if font.size:
                                    # 字体大小是Points，转换为像素（1pt = 1.33px at 96 DPI）
                                    # 对于150 DPI，1pt ≈ 2px
                                    font_size = int(font.size.pt * 2)
                    except:
                        pass
                    
                    text_items.append({
                        "text": text,
                        "position": {
                            "x": round(x, 2),
                            "y": round(y, 2),
                            "width": round(width, 2),
                            "height": round(height, 2),
                        },
                        "font_size": font_size
                    })
                except Exception as e:
                    # 如果提取位置失败，记录但继续
                    print(f"Warning: Failed to extract position for text '{text[:20]}...': {e}")
                    continue
        
        # 返回文本位置和尺寸信息
        return {
            "text_items": text_items,
            "slide_dimensions": {
                "width_inches": round(slide_width_inches, 2),
                "height_inches": round(slide_height_inches, 2),
                "width_pixels": round(slide_width_px, 2),  # 基于150 DPI
                "height_pixels": round(slide_height_px, 2),
                "dpi": resolution,
                "width_emu": slide_width_emu,
                "height_emu": slide_height_emu
            }
        }

