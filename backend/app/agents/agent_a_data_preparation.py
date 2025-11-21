# -*- coding: utf-8 -*-
# ===========================================================
# 文件：backend/app/agents/agent_a_data_preparation.py
# 功能：Agent A - 数据准备与题库生成（含 debug 输出）
# ===========================================================

import os
import re
from typing import List
from datetime import datetime

from app.agents.shared_state import shared_state
from app.agents.database.question_bank_storage import save_question_bank
from app.agents.models.quiz_models import Question, QuestionBank


# -----------------------------------------------------------
# 文本提取
# -----------------------------------------------------------

def extract_text_from_file(file_path: str) -> str:
    """
    从上传的文档中提取文本内容
    支持 .pdf / .pptx / .docx / .txt
    """
    ext = os.path.splitext(file_path)[1].lower()
    text = ""

    try:
        if ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        elif ext == ".docx":
            from docx import Document
            doc = Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])

        elif ext == ".txt":
            # ✅ 新增：纯文本文件支持
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        else:
            print(f"[⚠️ 不支持的文件类型] {ext}")

    except Exception as e:
        print(f"[⚠️ 文件解析失败] {e}")

    return text


# -----------------------------------------------------------
# 从文本中解析题目（增强版，支持多种格式）
# -----------------------------------------------------------

def parse_questions_from_text(raw_text: str) -> List[Question]:
    """
    从文本中解析题目，支持多种常见格式：
    1. 标准格式：1. 题干\n答案：xxx
    2. 选择题格式：1. 题干\nA. xxx\nB. xxx\nC. xxx\nD. xxx
    3. 简单编号：1. 题干（无答案）
    4. 问答格式：Question 1: xxx\nAnswer: xxx
    5. 中文格式：一、题干\n二、题干
    """
    questions = []
    if not raw_text.strip():
        print("👉 Debug: 输入文本为空，跳过正则解析。")
        return questions

    # 模式1: 标准数字编号 + 可选答案 (1. 题干 答案：xxx)
    pattern1 = r"(\d+)[\.\、\s]+(.+?)(?:答案[:：]\s*(.*?))?(?=\n\d+[\.\、\s]+|\n[一二三四五六七八九十][\.\、]|$)"
    matches1 = re.findall(pattern1, raw_text, re.S | re.I)
    
    # 模式2: 选择题格式 (1. 题干\nA. xxx\nB. xxx...)
    pattern2 = r"(\d+)[\.\、\s]+(.+?)(?:\n[A-D][\.\)]\s*.+){2,}(?:\n答案[:：]\s*([A-D]))?(?=\n\d+[\.\、]|\n[一二三四五六七八九十][\.\、]|$)"
    matches2 = re.findall(pattern2, raw_text, re.S | re.I)
    
    # 模式3: Question/Answer 格式
    pattern3 = r"(?:Question|问题)\s*(\d+)[:：\s]+(.+?)(?:(?:\n|\r\n)(?:Answer|答案)[:：\s]+(.+?))?(?=(?:\n|\r\n)(?:Question|问题)\s*\d+|$)"
    matches3 = re.findall(pattern3, raw_text, re.S | re.I)
    
    # 模式4: 中文数字编号（一、二、三...）
    pattern4 = r"([一二三四五六七八九十]+)[\.\、、]\s*(.+?)(?:答案[:：]\s*(.*?))?(?=\n[一二三四五六七八九十]+[\.\、、]|$)"
    matches4 = re.findall(pattern4, raw_text, re.S)
    
    # 合并所有匹配结果
    all_matches = []
    
    # 处理模式1的匹配
    for num, stem, ans in matches1:
        all_matches.append((num, stem.strip(), (ans or '').strip()))
    
    # 如果模式1没有匹配到，尝试其他模式
    if not all_matches:
        for num, stem, ans in matches2:
            all_matches.append((num, stem.strip(), (ans or '').strip()))
    
    if not all_matches:
        for num, stem, ans in matches3:
            all_matches.append((num, stem.strip(), (ans or '').strip()))
    
    if not all_matches:
        # 转换中文数字为阿拉伯数字
        chinese_to_num = {'一': '1', '二': '2', '三': '3', '四': '4', '五': '5',
                         '六': '6', '七': '7', '八': '8', '九': '9', '十': '10'}
        for cn_num, stem, ans in matches4:
            num = chinese_to_num.get(cn_num, str(len(all_matches) + 1))
            all_matches.append((num, stem.strip(), (ans or '').strip()))
    
    print(f"👉 Debug: 匹配到题目数量 = {len(all_matches)}")

    # 构建 Question 对象
    for idx, (num, stem, ans) in enumerate(all_matches, 1):
        # 清理题干（去除多余空白和换行）
        stem_clean = re.sub(r'\s+', ' ', stem).strip()
        
        # 检测题型
        question_type = "short_answer"
        if re.search(r'[A-D][\.\)]\s*.+', stem):
            question_type = "multiple_choice"
        elif any(keyword in stem_clean for keyword in ['编程', '代码', '算法', 'code', 'program']):
            question_type = "programming"
        elif any(keyword in stem_clean for keyword in ['判断', '对错', 'true', 'false']):
            question_type = "true_false"
        
        q = Question(
            id=f"Q{idx:03d}",
            stem=stem_clean[:1000] if len(stem_clean) > 1000 else stem_clean,  # 限制长度
            answer=(ans or '').strip() or "（待补充）",
            difficulty="medium",
            knowledge_points=[],
            question_type=question_type
        )
        questions.append(q)

    print(f"👉 Debug: parse_questions_from_text() 返回 {len(questions)} 题")
    return questions


# -----------------------------------------------------------
# Agent A 不再生成默认题库，如果无法提取题目则返回空列表
# 由后续的 Agent E 通过 LLM 基于文本内容生成题目
# -----------------------------------------------------------


# -----------------------------------------------------------
# 主执行函数
# -----------------------------------------------------------

def run_agent_a(conversation_id: str, file_paths: List[str]) -> QuestionBank:
    print(f"🧩 [Agent A] 开始生成题库，会话ID: {conversation_id}，共 {len(file_paths)} 个文件")

    # -------------------------------------------------------
    # 🔍 自动扫描当前会话上传目录中的 txt（仅当 file_paths 为空）
    # -------------------------------------------------------
    if (not file_paths) or (len(file_paths)==1 and file_paths[0] == "__AUTO__"):
        print(f"👉 Debug: file_paths 为空，正在扫描会话 [{conversation_id}] 的样本文件...")

        # 自动定位 backend 根目录（agent_a 所在的目录是 backend/app/agents/）
        BASE_DIR = os.path.dirname(os.path.dirname(os.path.dirname(__file__)))
        # ✅ 只扫描当前会话的目录
        base_dir = os.path.join(BASE_DIR, "uploads", "exercises", conversation_id, "samples")
        detected_files = []

        if os.path.exists(base_dir):
            # 遍历 exercises/<conversation_id>/samples/<folder>/text.txt
            for root, dirs, files in os.walk(base_dir):
                for f in files:
                    # 优先使用 text.txt（解析后的文本），也支持其他 .txt 文件
                    if f.lower() == "text.txt" or f.lower().endswith(".txt"):
                        full = os.path.join(root, f)
                        detected_files.append(full)

        if not detected_files:
            print(f"⚠️ 未在会话 [{conversation_id}] 的样本目录下找到任何 txt 文件，将进入默认题库模式。")
        else:
            print(f"👉 Debug: 在会话 [{conversation_id}] 中找到 {len(detected_files)} 个样例文件：")
            for df in detected_files:
                print("   -", df)

            file_paths = detected_files


    all_text = ""

    if file_paths:
        for fp in file_paths:
            if os.path.exists(fp):
                print(f"👉 Debug: 正在解析文件 {fp}")
                extracted = extract_text_from_file(fp)
                if extracted:
                    print(f"👉 Debug: 从文件中提取到 {len(extracted)} 字符文本。")
                    all_text += extracted + "\n"
                else:
                    print(f"[⚠️ 文件无可提取文本] {fp}")
            else:
                print(f"[⚠️ 文件不存在] {fp}")
    else:
        print("ℹ️ 未检测到上传的样例试卷，将进入自动生成模式。")

    print(f"👉 Debug: 提取文本长度 = {len(all_text)}")

    questions = parse_questions_from_text(all_text)
    print(f"👉 Debug: 正则解析后题目数 = {len(questions)}")

    # ⚠️ 如果正则解析失败，尝试智能分段处理
    if not questions and all_text.strip():
        print("⚠️ 正则解析未找到题目，尝试智能分段处理...")
        questions = create_questions_from_content(all_text)
        print(f"👉 智能分段生成 {len(questions)} 个题目框架")

    # 如果仍然没有题目，返回空题库交由 Agent E 处理
    if not questions:
        print("⚠️ 未从样本文件中解析到任何题目，将返回空题库交由 Agent E 处理")
        print(f"   提取的文本前500字符: {all_text[:500]}")
    else:
        # 保存原始文本到 shared_state，供 Agent E 使用
        shared_state.source_text = all_text

    qb = QuestionBank(questions=questions, source_docs=file_paths)
    shared_state.question_bank = qb

    save_path = save_question_bank(conversation_id, qb)
    print(f"✅ 题库已生成并保存到: {save_path} （共 {len(questions)} 题）")

    print("👉 Debug: run_agent_a() 执行完毕。")
    return qb


# -----------------------------------------------------------
# 智能分段处理：从内容中创建题目框架
# -----------------------------------------------------------

def create_questions_from_content(text: str) -> List[Question]:
    """
    当无法提取标准题目时，将文本智能分段，
    为每个段落创建题目框架，供 Agent E 基于这些段落生成题目。
    """
    questions = []
    
    # 按段落分割（双换行或多个换行）
    paragraphs = re.split(r'\n\s*\n+', text.strip())
    
    # 过滤太短的段落（少于20个字符）
    meaningful_paragraphs = [p.strip() for p in paragraphs if len(p.strip()) > 20]
    
    if not meaningful_paragraphs:
        print("   智能分段：未找到有意义的段落")
        return []
    
    print(f"   智能分段：找到 {len(meaningful_paragraphs)} 个段落")
    
    # 为每个段落创建题目框架（最多取前10个段落）
    for idx, para in enumerate(meaningful_paragraphs[:10], 1):
        # 提取段落关键词（取前50个字作为题干提示）
        stem_hint = para[:100].replace('\n', ' ').strip()
        
        q = Question(
            id=f"Q{idx:03d}",
            stem=f"基于以下内容出题：{stem_hint}...",
            answer="（待 Agent E 生成）",
            difficulty="medium",
            knowledge_points=[],
            question_type="short_answer"
        )
        questions.append(q)
    
    return questions
